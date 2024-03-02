import streamlit as st
import time
import random
import os
import re
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from docx import Document
from PIL import Image
import fitz
import pptx
import docx
import io
from pptx import Presentation
from pylatexenc import latex2text


if "user_chat" not in st.session_state:
    st.session_state.user_chat = []
if "ai_chat" not in st.session_state:
    st.session_state.ai_chat = []

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

st.set_page_config(
    page_title="Docs Sherlock",
    page_icon="🔍",
    layout="wide"
)
st.title("Docs Sherlock 🔍")


def stream_data(test):
    for word in test.split(' '):
        yield word + " "
        time.sleep(0.09)


def extract_images(file, extension):
    images = []

    if extension == "pdf":
        pdf = fitz.open(file)

        for page_number in range(len(pdf)):
            page = pdf.load_page(page_number)
            Objects = page.get_images(full=True)

            for _, img_info in enumerate(Objects):
                base_image = pdf.extract_image(img_info[0])
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)


    elif extension == "pptx" or extension == "ppt":
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = Image.open(io.BytesIO(shape.image.blob))
                    images.append(image)


    elif extension == "docx" or extension == "doc":
        doc = docx.Document(file)
        
        rel_parts = doc.part.rels.values()
        
        for rel in rel_parts:
            if "image" in rel.reltype:
                image_data = rel.target_part.blob 
                image = Image.open(io.BytesIO(image_data))
                images.append(image)

    return images


def image_description(images_extracted):
    model = genai.GenerativeModel('gemini-1.0-pro-vision-latest')
    return_str = ""
    for image in images_extracted:
        try:
            response = model.generate_content(image)
            summary = response.text

            return_str = return_str + f"Summary of Fig. {i + 1}: {summary}\n\n"

            i += 1
        except:
            pass
    
    return return_str


def pdf_docs_ppt_image(file, extension):
    images_extracted = extract_images(file, extension)

    if len(images_extracted) == 0:
        return -1

    model = genai.GenerativeModel('gemini-1.0-pro-vision-latest')

    return_str = ""

    i = 0 

    for image in images_extracted:
        try:
            response = model.generate_content(image)
            summary = response.text

            return_str = return_str + f"Summary for Fig. {i + 1}: {summary}\n\n"

            i += 1
        except:
            pass
    
    
    return return_str


def download_log():
    str = ""
    for user_text, ai_text in zip(st.session_state.user_chat, st.session_state.ai_chat):
        str = str + f"User: {user_text}" + f"\nPDF Sherlock: {ai_text}\n"
    return str


def space_remove(str):
    return re.sub(r'\s+', ' ', str).strip()


def latex_to_txt(tex_file):
    latex_file = tex_file

    with open(latex_file, "r") as f:
        latex_content = f.read()

    text_content = latex2text.latex2text(latex_content)

    image_pattern = r"\\includegraphics\[.*\]\{(.*?)\}"

    image_names = re.findall(image_pattern, latex_content)

    formated_string = space_remove(text_content)

    return formated_string, image_names


def get_doc_text(docx_file):
    doc = Document(docx_file)
    full_text = []
    for paragraph in doc.paragraphs:
        full_text.append(paragraph.text)

    return_str = '\n'.join(full_text)
    return_str = return_str.lower()
    return return_str


def get_pdf_text(pdfs):
    text = ""
    for pdf in pdfs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def get_txt_text(text_docs):
    text = text_docs.getvalue().decode("utf-8")
    return text


def pdf_txt_decide(mode, docs):
    if mode == "pdf":
        text = get_pdf_text([docs])
    elif mode == "doc" or mode == "docx":
        text = get_doc_text(docs)
    elif mode == "ppt" or mode == "pptx":
        text = get_ppt_text(docs)
    else:
        text = get_txt_text(docs)
    return text.lower()


def get_ppt_text(ppt_docs):
    prs = Presentation(ppt_docs)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = text + str(shape.text) + "\n"

    return text


def text_to_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)

    return chunks


def docs_to_chunks(docs):
    r_splitter = RecursiveCharacterTextSplitter(
    separators = ["\n\n", "\n", " "],
    chunk_size =100000 ,
    chunk_overlap  = 1000, 
    )
    chunks = r_splitter.split_documents(docs)

    return chunks


def chunk_to_FAISS(chunks, mode):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    if mode == 0:
        vector_store = FAISS.from_texts(chunks, embedding=embeddings)
        vector_store.save_local("faiss_model")
    elif mode == 1:
        vector_store = FAISS.from_documents(chunks, embedding=embeddings)
        vector_store.save_local("faiss_model")


def ai_generative_reply():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-1.0-pro",
                             temperature=0.3)

    prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
    reply = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return reply


def answer(question):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")

    new_db = FAISS.load_local("faiss_model", embeddings)
    docs = new_db.similarity_search(question)

    reply = ai_generative_reply()

    response = reply(
    {"input_documents":docs, "question": question}
    , return_only_outputs=True)

    return str(response['output_text'])

    
side = st.sidebar
side.title("Upload your document")

uploaded_file = side.file_uploader("DOC/DOCX/PPT/PPTX/PDF/TXT/.TEX", accept_multiple_files=False, type=['doc','docx','ppt', 'pptx','pdf', 'txt', '.tex'])
if side.button("Submit"):
    if uploaded_file is not None:
        file_extension = uploaded_file.name.split('.')[-1]

        st.session_state.user_chat = []
        st.session_state.ai_chat = []

        if file_extension == 'tex':
            tex_file = uploaded_file.name
            text, images = latex_to_txt(tex_file)
            chunks = text_to_chunks(text)
            chunk_to_FAISS(chunks, 0)
            st.success("Done")

        if file_extension in ['pdf', 'txt', 'docx', 'doc','ppt', 'pptx']:
            with st.spinner("Processing Data..."):
                text = pdf_txt_decide(file_extension,uploaded_file)
                chunks = text_to_chunks(text)

                if file_extension in ['pdf', 'docx', 'doc','ppt', 'pptx']:
                    
                    images_summary = pdf_docs_ppt_image(str(uploaded_file.name), file_extension)
                    if images_summary != -1:
                        image_chunks = text_to_chunks(images_summary)
                        chunks.extend(image_chunks)

                chunk_to_FAISS(chunks, 0)
                st.success("Done")


io_container= st.container(height=500)

status_messages = [
    "Your request is being processed. Please wait.",
    "Processing your request. Hang tight!",
    "Currently processing your data. Almost there!",
    "We're working on it! Your text is being processed.",
    "Sit tight! Processing in progress.",
    "Text is in the queue for processing.",
    "Processing your text. Thanks for your patience!",
    "Our backend is crunching the data. Please be patient.",
    "Text processing in progress. It won't be long now!",
    "Processing underway. We appreciate your patience."
    ]


if prompt:= io_container.chat_input("Text Sherlock..."):
    random_element = random.choice(status_messages)

    status = io_container.empty()
    status.text(random_element)
    
    ai_answer = answer(prompt.lower())

    status.empty()

    io_container.chat_message("user").write(prompt)
    st.session_state.user_chat.append(prompt)

    expander = io_container.expander("🤖 Docs Sherlock: ", expanded=True)
    expander.write_stream(stream_data(ai_answer))
    st.session_state.ai_chat.append(ai_answer)

with st.expander("Chat Log"):
    text_data = download_log().encode("utf-8'")
    st.download_button('Download Log', data=text_data, mime='text/plain', file_name="log.txt")
        
    for user_text, ai_text in zip(st.session_state.user_chat, st.session_state.ai_chat):
        st.chat_message("user").write(f"User: {user_text}")
        st.chat_message("assistant").write(f"Docs Sherlock: {ai_text}")
